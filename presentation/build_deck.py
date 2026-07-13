"""
Generates the OpenGov "Sr Data Platform Engineer" case-study presentation (Part 1:
Platform Architecture & Data Foundation) as a native, editable .pptx.

Run:  python presentation/build_deck.py
Out:  presentation/OpenGov_Data_Platform_CaseStudy.pptx

Design: 16:9, navy theme, white-circle/navy-arrow OpenGov logo mark, one big
data-flow diagram (slide 3), an RBAC diagram (slide 10), a decision-matrix table
(slide 4), and speaker notes on every slide to drive a ~20 minute talk.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, "image.png")   # OpenGov mark

# ----------------------------------------------------------------------------- theme (OpenGov indigo brand)
NAVY   = RGBColor(0x19, 0x15, 0x52)   # deep indigo (title bars, dark slides, dark boxes)
NAVY2  = RGBColor(0x2A, 0x24, 0x70)   # secondary deep indigo
BLUE   = RGBColor(0x3B, 0x34, 0xB8)   # rich indigo (primary box fills)
ACCENT = RGBColor(0x4A, 0x3F, 0xFF)   # OpenGov brand indigo (rules, bullets, arrows)
LAV    = RGBColor(0xA7, 0x9F, 0xFF)   # light lavender (accent text on dark)
LIGHT  = RGBColor(0xEC, 0xEB, 0xFB)   # light panel
CARD   = RGBColor(0xF5, 0xF5, 0xFC)   # card fill
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1B, 0x1B, 0x33)   # near-black body text
GRAY   = RGBColor(0x5C, 0x5B, 0x77)   # muted
GRAYINK= RGBColor(0x3E, 0x3D, 0x5C)   # sub-bullet text
GREEN  = RGBColor(0x2E, 0xA0, 0x74)   # "build" accent
AMBER  = RGBColor(0xE0, 0x8A, 0x2B)   # "buy" accent
CODEFG = RGBColor(0xE6, 0xE9, 0xF5)   # code text
CODECMT= RGBColor(0x8E, 0xD6, 0xB0)   # code comment

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------- helpers
def blank():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _noline(shape):
    shape.line.fill.background()

def rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, line_w=1.0):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        _noline(s)
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s

def text_in(shape, text, size=14, color=INK, bold=False, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, font=FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = _dash(text)
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tf

def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tb, tf

def place_logo(slide, cx, cy, size):
    """Place the OpenGov logo image centered at (cx, cy), clipped to a circle."""
    pic = slide.shapes.add_picture(LOGO, Inches(cx - size/2), Inches(cy - size/2),
                                   Inches(size), Inches(size))
    # clip the (square) image into a circle via preset ellipse geometry
    spPr = pic._element.spPr
    old = spPr.find(qn('a:prstGeom'))
    if old is not None:
        spPr.remove(old)
    geom = spPr.makeelement(qn('a:prstGeom'), {'prst': 'ellipse'})
    geom.append(geom.makeelement(qn('a:avLst'), {}))
    spPr.append(geom)
    return pic

def title_bar(slide, title, kicker=None):
    bar = rect(slide, 0, 0, SW.inches, 1.12, NAVY)
    rect(slide, 0, 1.12, SW.inches, 0.055, ACCENT)          # accent rule
    tb, tf = textbox(slide, 0.55, 0.12, 10.8, 0.92)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kicker:
        p0 = tf.paragraphs[0]; r0 = p0.add_run(); r0.text = kicker.upper()
        r0.font.size = Pt(11); r0.font.bold = True; r0.font.color.rgb = LAV; r0.font.name = FONT
        p0.space_after = Pt(1)
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    r = p.add_run(); r.text = tidy(title)
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    place_logo(slide, SW.inches - 0.72, 0.56, 0.6)          # mark in the bar

PAGE = [1]
def footer(slide, n=None):
    PAGE[0] += 1
    tb, tf = textbox(slide, 0.55, 7.06, 8.0, 0.35)
    p = tf.paragraphs[0]; r = p.add_run()
    r.text = "OpenGov Data Platform  ·  Case Study  ·  Akash Pahilwan"
    r.font.size = Pt(9); r.font.color.rgb = GRAY; r.font.name = FONT
    tb2, tf2 = textbox(slide, SW.inches - 1.4, 7.06, 0.9, 0.35)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT; r2 = p2.add_run()
    r2.text = str(PAGE[0]); r2.font.size = Pt(9); r2.font.color.rgb = GRAY; r2.font.name = FONT

# lowercase technical tokens that must NOT be sentence-cased
_TECH = {"dbt", "event_id", "stage_name", "is_deleted", "account_id", "opportunity_id",
         "payload", "sync_config.py", "apply_pii_tags.py", "ingest_page_views.py"}

def _dash(t):
    """Replace em/en dashes with commas (avoid the em-dash 'AI tell')."""
    if not t:
        return t
    return (t.replace(" — ", ", ").replace(" – ", ", ")
             .replace("—", ", ").replace("–", ", "))

def tidy(t):
    """Dash-clean + sentence-case the first word (unless it's a technical token)."""
    t = _dash(t)
    if not t:
        return t
    t = t.lstrip()
    while t[:1] in (",", ";"):          # a body that began with a dash now leads with ", "
        t = t[1:].lstrip()
    i = 0
    while i < len(t) and not t[i].isalpha():
        i += 1
    if i < len(t):
        w = t[i:].split(" ", 1)[0].strip(".,:;/()")
        base = w.lower().split("-")[0].split(".")[0]   # e.g. "dbt-owned" -> "dbt"
        if w.lower() not in _TECH and base not in _TECH and "_" not in w:
            t = t[:i] + t[i].upper() + t[i + 1:]
    return t

def bullets(slide, items, x, y, w, h, base=18):
    """items: list of (level:int, lead:str|None, body:str)."""
    tb, tf = textbox(slide, x, y, w, h)
    for i, (level, lead, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8 if level == 0 else 3)
        p.line_spacing = 1.07
        size = base if level == 0 else base - 3
        g = p.add_run()
        g.text = "▸  " if level == 0 else "        ·  "
        g.font.size = Pt(size); g.font.name = FONT
        g.font.color.rgb = ACCENT if level == 0 else GRAY
        if lead:
            r = p.add_run(); r.text = _dash(lead) + ("   " if body else "")
            r.font.bold = True; r.font.size = Pt(size); r.font.color.rgb = INK; r.font.name = FONT
        if body:
            r2 = p.add_run(); r2.text = tidy(body)
            r2.font.size = Pt(size); r2.font.name = FONT
            r2.font.color.rgb = INK if level == 0 else GRAYINK
    return tb

def takeaway(slide, text, y=6.35):
    bar = rect(slide, 0.55, y, 12.23, 0.5, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, 0.55, y, 0.09, 0.5, ACCENT)
    tb, tf = textbox(slide, 0.8, y, 11.8, 0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r0 = p.add_run(); r0.text = "Takeaway   "; r0.font.bold = True
    r0.font.size = Pt(13); r0.font.color.rgb = BLUE; r0.font.name = FONT
    r1 = p.add_run(); r1.text = tidy(text)
    r1.font.size = Pt(13); r1.font.color.rgb = INK; r1.font.name = FONT

def code_box(slide, x, y, w, h, lines, title=None, fs=10):
    """Dark code sketch with light text and green comments."""
    rect(slide, x, y, w, h, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb, tf = textbox(slide, x + 0.16, y + 0.12, w - 0.32, h - 0.24)
    tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    if title:
        p = tf.paragraphs[0]; p.space_after = Pt(6)
        r = p.add_run(); r.text = title
        r.font.name = "Consolas"; r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = LAV
        first = False
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.0; p.space_after = Pt(1)
        if "--" in ln:
            idx = ln.index("--"); code, comment = ln[:idx], ln[idx:]
        else:
            code, comment = ln, None
        if code:
            r1 = p.add_run(); r1.text = code
            r1.font.name = "Consolas"; r1.font.size = Pt(fs); r1.font.color.rgb = CODEFG
        if comment:
            r2 = p.add_run(); r2.text = comment
            r2.font.name = "Consolas"; r2.font.size = Pt(fs); r2.font.color.rgb = CODECMT

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = _dash(text)

def content(title, kicker=None):
    s = blank(); set_bg(s, WHITE); title_bar(s, title, kicker); return s

# ============================================================================= SLIDE 1 — Title
s = blank(); set_bg(s, NAVY)
rect(s, 0, 0, 0.28, SH.inches, ACCENT)                    # left accent band
place_logo(s, 2.0, 2.35, 1.35)
tb, tf = textbox(s, 1.9, 3.45, 10.8, 2.6)
p = tf.paragraphs[0]; r = p.add_run(); r.text = "OpenGov Data Platform"
r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "Architecture & Data Foundation  ·  Hands-On Build (RevOps)"
r2.font.size = Pt(20); r2.font.color.rgb = LAV; r2.font.name = FONT
p3 = tf.add_paragraph(); p3.space_before = Pt(18)
r3 = p3.add_run(); r3.text = "Case Study Presentation"
r3.font.size = Pt(15); r3.font.color.rgb = RGBColor(0xC7,0xD3,0xE0); r3.font.name = FONT
p4 = tf.add_paragraph()
r4 = p4.add_run(); r4.text = "Akash Pahilwan  ·  Senior Data Platform Engineer"
r4.font.size = Pt(15); r4.font.bold = True; r4.font.color.rgb = WHITE; r4.font.name = FONT
notes(s, "Open with the thesis: I build the factory — the governed foundation that "
         "lets every domain team ship analytics and AI on top, safely and fast. Today "
         "I'll walk foundation -> pipelines -> self-service -> trust & governance, then "
         "close with a forward-looking bet on making the platform AI-ready.")

# ============================================================================= SLIDE 2 — Thesis
s = content("The platform is a product — I build the factory", "Foundation")
bullets(s, [
    (0, "The mandate:", "one governed foundation powering analytics and AI across GTM, Finance, Product, and HR/People."),
    (0, "Federated, not centralized:", "domain teams build on a paved path with guardrails — self-service, not a ticket queue."),
    (0, "Four principles guide every decision:", ""),
    (1, "Governed by default", "— RBAC, masking, and lineage are built in, not bolted on."),
    (1, "Self-service paved path", "— golden templates + IaC modules; onboard a domain in a PR."),
    (1, "Observable & SLA-driven", "— freshness, quality, and cost measured from day one."),
    (1, "IaC-first & reproducible", "— everything is code, tested in CI, rebuildable in any environment."),
], 0.7, 1.5, 12.0, 4.6, base=19)
takeaway(s, "Trust and speed are not a trade-off — the right guardrails deliver both.")
footer(s, 2)
notes(s, "Frame the role: a platform engineer's customer is other engineers/analysts. "
         "The job is to make the right thing the easy thing. Stress federation with "
         "guardrails — central team owns the paved path, domains own their data products. "
         "These four principles recur on every later slide.")

# ============================================================================= SLIDE 3 — Architecture diagram
s = content("Target-state architecture: source → dashboard", "Foundation")

def dbox(x, y, w, h, txt, fill, tc=WHITE, fs=11, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    b = rect(s, x, y, w, h, fill, shape=shape)
    text_in(b, txt, size=fs, color=tc, bold=bold)
    return b

def arrow(x, y, w=0.35, h=0.3):
    a = rect(s, x, y, w, h, ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)
    return a

def col_label(x, y, w, txt):
    tb, tf = textbox(s, x, y, w, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run()
    r.text = txt.upper(); r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = BLUE; r.font.name = FONT

top = 1.75
# Col 1 — Sources
col_label(0.55, top-0.32, 2.25, "Sources")
srcs = [("Salesforce  (GTM)", BLUE), ("Marketo  (GTM)", BLUE), ("NetSuite  (Finance)", BLUE),
        ("Telemetry  (Product)", NAVY2), ("Workday  (HR)", BLUE)]
for i,(t,c) in enumerate(srcs):
    dbox(0.55, top + i*0.72, 2.25, 0.58, t, c, fs=10)

# Col 2 — Ingestion
col_label(3.15, top-0.32, 2.5, "Ingestion")
dbox(3.15, top+0.15, 2.5, 1.05, "Fivetran\nmanaged connectors (SaaS)", AMBER, fs=11, bold=True)
dbox(3.15, top+1.55, 2.5, 1.35, "Custom Python\nAWS S3 + Lambda\n(telemetry & APIs)", GREEN, fs=11, bold=True)

# Col 3 — Snowflake
col_label(6.05, top-0.32, 3.65, "Snowflake  +  dbt")
snow = rect(s, 6.05, top, 3.65, 3.55, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=BLUE, line_w=1.0)
layers = [("RAW: source-system schemas", NAVY), ("STAGING: cleaned / conformed", BLUE),
          ("MARTS: domain, business-ready", BLUE), ("SANDBOX: analyst scratch", GRAY)]
for i,(t,c) in enumerate(layers):
    dbox(6.28, top+0.24 + i*0.72, 3.2, 0.58, t, c, fs=10)
tb, tf = textbox(s, 6.28, top+3.02, 3.2, 0.4)
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; r=p.add_run()
r.text="dbt transforms & tests RAW→STAGING→MARTS"; r.font.size=Pt(9); r.font.italic=True
r.font.color.rgb=BLUE; r.font.name=FONT

# Col 4 — Consumers
col_label(10.05, top-0.32, 2.7, "Consumers")
cons = [("BI & Dashboards", BLUE), ("AI / ML  ·  LLM + RAG", NAVY2), ("Reverse ETL → apps", BLUE)]
for i,(t,c) in enumerate(cons):
    dbox(10.05, top+0.35 + i*0.95, 2.7, 0.72, t, c, fs=11, bold=(i==1))

# arrows between columns
arrow(2.83, top+1.4); arrow(5.68, top+1.4); arrow(9.73, top+1.4)

# cross-cutting foundation band
band = rect(s, 0.55, 5.55, 12.18, 0.62, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text_in(band, "FOUNDATION   ·   Terraform (IaC)   ·   RBAC + Dynamic Masking   ·   Observability & Data SLAs   ·   CI/CD",
        size=12, color=WHITE, bold=True)
footer(s, 3)
notes(s, "This is the anchor diagram — walk it left to right. Sources land via Fivetran "
         "(connectorized SaaS) or a custom Python path on AWS (S3 + Lambda) for telemetry "
         "and bespoke APIs. Everything lands in Snowflake RAW, then dbt promotes RAW→STAGING"
         "→MARTS, with SANDBOX for exploration. Consumers are BI, AI/ML + LLM/RAG, and "
         "reverse ETL. The navy band is the cross-cutting foundation that every layer sits "
         "on: IaC, RBAC+masking, observability, CI/CD. Invite the 'why Fivetran here, Lambda "
         "there' question — next slide answers it.")

# ============================================================================= SLIDE 4 — Fivetran vs custom (table)
s = content("Ingestion: buy vs. build", "Pipelines")
rows = [
    ("Dimension", "Fivetran  (buy)", "Custom Python on AWS  (build)"),
    ("Best for", "Connectorized SaaS sources", "High-volume / bespoke sources, no connector"),
    ("Examples", "Salesforce, NetSuite, Workday, Marketo", "Product telemetry, operational APIs"),
    ("Schema drift", "Schema-on-write — auto-evolves its tables", "Schema-on-read — VARIANT + contracts"),
    ("Latency", "Scheduled batch", "Near-real-time — S3 event → Lambda"),
    ("Control", "Low (fully managed)", "Full — custom logic & cost tuning"),
    ("Cost model", "Per-MAR (rows synced)", "Compute — S3 / Lambda / warehouse"),
]
nrows, ncols = len(rows), 3
gf = s.shapes.add_table(nrows, ncols, Inches(0.6), Inches(1.5), Inches(12.13), Inches(4.2))
tbl = gf.table
tbl.columns[0].width = Inches(2.3); tbl.columns[1].width = Inches(4.6); tbl.columns[2].width = Inches(5.23)
for ci in range(ncols):
    tbl.rows[0].height = Inches(0.62)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = _dash(val)   # table cells bypass tidy(); strip em-dashes here too
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.08)
        para = cell.text_frame.paragraphs[0]
        run = para.runs[0]; run.font.name = FONT
        if ri == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = WHITE
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else CARD
            run.font.size = Pt(12.5); run.font.color.rgb = INK
            if ci == 0:
                run.font.bold = True; run.font.color.rgb = BLUE
            if ci == 1 and ri > 0:
                run.font.color.rgb = AMBER if val else INK
            if ci == 2 and ri > 0:
                run.font.color.rgb = GREEN
takeaway(s, "Buy the commodity (SaaS connectors), build the differentiator (telemetry scale, cost, custom logic).")
footer(s, 4)
notes(s, "The decision rule: buy vs build on control, cost curve, and whether a reliable "
         "connector exists. Fivetran wins for SaaS — someone else maintains the connector "
         "and absorbs API changes. Telemetry is the opposite: huge volume, simple shape, "
         "cost at Fivetran's per-row pricing would be brutal, and we want near-real-time "
         "via S3 events → Lambda. So we own that path. Be ready to defend the cost math.")

# (SLIDE removed — generic connector theory; API pagination/rate-limiting fielded
#  verbally, and our ingestion is shown concretely on "Reliable landing" + Part 2 P2.4)

# ============================================================================= SLIDE 5B — Schema drift (CI/CD)
s = content("Schema drift when objects ship via CI/CD", "Pipelines")
bullets(s, [
    (0, "Two schemas, two rules:", "the upstream source shape is uncontrolled; our deployed objects change only through reviewed PRs."),
    (0, "Custom path — schema-on-read:", "RAW lands each record as a VARIANT; a new source field needs zero DDL, zero deploy — typed at read time in dbt."),
    (0, "Detect, don't auto-alter:", "compare the payload to a registered contract — additive → log + surface; breaking → alert + quarantine. CI/CD never mutates a live table from drift."),
    (0, "Promote a field = a pull request:", "expose it by adding one typed column in dbt; CI runs contracts + tests; review + merge deploys it."),
    (0, "Fivetran — schema-on-write:", "it types columns and ALTERs its own RAW tables on drift; the dbt source contract + tests are our CI tripwire one layer down(If needed)."),
], 0.65, 1.5, 7.05, 4.7, base=16.5)
code_box(s, 8.0, 1.6, 4.78, 4.55, [
    "-- RAW - schema-on-read (stable DDL)",
    "CREATE OR ALTER TABLE",
    "  raw.product_events.page_views (",
    "  event_id  STRING,   -- promoted key",
    "  loaded_at TIMESTAMP_NTZ,",
    "  _filename STRING,",
    "  payload   VARIANT   -- new fields",
    ");                    -- land here, no PR",
    "",
    "-- STAGING (dbt) - type late",
    "SELECT",
    " payload:event_id::string  AS event_id,",
    " payload:properties.duration_ms::int",
    "                           AS duration_ms",
    "FROM {{ source('product_events',",
    "               'page_views') }}",
    "-- expose a field = one line, via PR",
], title="drift: land loose, type in a PR")
takeaway(s, "Upstream drift lands safely in VARIANT; the deployed schema changes only through a reviewed, tested PR — drift is a blocking CI signal, not silent breakage.")
footer(s)
notes(s, "This directly answers 'we deploy objects via CI/CD — how does drift fit?'. Separate "
         "the two schemas. The SOURCE shape is uncontrolled and can change any time; OUR "
         "objects are code and change only via a reviewed, tested PR. RAW uses a VARIANT "
         "column so upstream additions need no DDL and never break a load — CI/CD does not "
         "run on every upstream change. Drift is detected against a contract and surfaced "
         "(additive) or quarantined + alerted (breaking), but never auto-alters a deployed "
         "table. To make a new field real, an engineer adds one typed column in dbt via a "
         "PR; contracts + tests run in CI as the tripwire. Fivetran is the contrast: it is "
         "schema-on-WRITE — it types columns and ALTERs its own RAW tables when the source "
         "drifts — so the read-time boundary moves one layer down: dbt sources select only "
         "the columns we depend on, and the source contract + tests fail CI before bad data "
         "flows downstream. Crisp framing if pushed: schema-on-read is a property of the RAW "
         "layer. In the custom path we own RAW, so we land VARIANT and type late; in the "
         "Fivetran path Fivetran owns and evolves RAW, so our read-time contract is the dbt "
         "source. Either way, a new field enters the deployed model only through a reviewed "
         "PR. Use CREATE OR ALTER / additive, idempotent migrations so redeploys are safe "
         "on half-failures.")

# ============================================================================= SLIDE 6 — Reliable landing
s = content("Reliable landing in Snowflake", "Pipelines")
bullets(s, [
    (0, "Partitioned landing zone:", "S3 keys as source/table/ingest_date=… — immutable, append-only, replayable."),
    (0, "File formats:", "JSON → VARIANT for semi-structured; Parquet for large columnar; compressed; bounded file sizes."),
    (0, "Load path:", "external stage + COPY INTO, or Snowpipe on arrival; capture METADATA$FILENAME for lineage."),
    (0, "Quality gates at load:", "validate required fields; bad rows → _QUARANTINE table; never fail the whole batch on one row."),
    (0, "Observability:", "a _LOAD_LOG row per file — records processed, quarantined, load timestamp."),
    (0, "Idempotent:", "dedup on natural key so re-loading a file is a no-op."),
], 0.7, 1.55, 12.0, 4.4, base=18)
takeaway(s, "RAW is immutable truth — validate and quarantine at the door, log every load.")
footer(s, 6)
notes(s, "RAW is the system of record — never mutate it. Partition for replay and cheap "
         "pruning. The quarantine + load-log pattern gives operability: you can always "
         "answer 'did the 09:00 file land, and how many rows did we reject?' This maps "
         "one-to-one to the PAGE_VIEWS ingestion I build in Part 2, including the "
         "intentional null user_id record that gets quarantined.")

# ============================================================================= SLIDE 7 — Layer model
s = content("Snowflake layer model & environment strategy", "Foundation")
bullets(s, [
    (0, "Database-per-layer:", "RAW → STAGING → MARTS → SANDBOX(Ideal); clear ownership and blast-radius per layer."),
    (1, "RAW", "— loader-owned, immutable, source-system schemas, no business logic."),
    (1, "STAGING", "— dbt-owned; typed, renamed, conformed; light transforms."),
    (1, "MARTS", "— dbt-owned; business-ready, domain schemas, tested & documented."),
    (1, "SANDBOX", "— analyst-writable scratch; readers & developers read-only, admin full; safe to break, no downstream deps."),
    (0, "Warehouse strategy:", "separate ingest / transform / BI warehouses; right-sized; auto-suspend for cost."),
    (0, "Environments:", "dev and prod isolated; promotion through CI, never manual edits in prod."),
], 0.7, 1.5, 12.0, 4.7, base=18)
takeaway(s, "Layers give ownership and blast-radius control; warehouses give workload & cost isolation.")
footer(s, 7)
notes(s, "Database-per-layer (vs schema-per-layer) gives cleaner grants and isolation, and "
         "scales to many domains: RAW holds source schemas, MARTS holds domain schemas. "
         "Separate warehouses stop a heavy dbt run from starving BI, and make cost "
         "attributable. Prod is only ever changed through CI.")

# ============================================================================= SLIDE 8 — dbt multi-domain
s = content("dbt for many domains — without a monolith", "Pipelines")
bullets(s, [
    (0, "Layered models:", "staging (1:1 with source) → intermediate → marts (by domain)."),
    (0, "Organize by domain, not by one giant project:", "folders + CODEOWNERS per domain (revops, finance, product, hr)."),
    (0, "dbt Mesh:", "each domain exposes PUBLIC models via contracts; consumers depend on stable interfaces, not internals."),
    (0, "Quality is code:", "generic + singular tests, model contracts, source freshness(ideal), exposures for BI/AI lineage."),
    (0, "Incremental vs full-refresh:", "incremental for high-volume events (telemetry); full-refresh for small dimensions."),
    (0, "Fast CI(Ideal):", "slim CI builds only state:modified+ — seconds, not full rebuilds."),
], 0.7, 1.55, 12.0, 4.5, base=18)
takeaway(s, "Contracts + dbt Mesh let domains move independently while consumers stay safe.")
footer(s, 8)
notes(s, "The 1-to-10-domains question lives here. Answer: don't grow one monolith — split "
         "by domain with clear ownership, and use dbt Mesh contracts as the interface "
         "between domains. Public models are the API; internals can refactor freely. "
         "Incremental vs full-refresh is a volume/cost decision. Mention mart_revops__"
         "pipeline (Part 2) would be exposed as a PUBLIC contracted model.")

# ============================================================================= SLIDE 9 — Self-service onboarding
s = content("Self-service: onboard a domain in a pull request", "Self-service")
bullets(s, [
    (0, "One Terraform module — domain_workspace:", "input a domain name, get a full, governed workspace."),
    (1, "Provisions", "schemas (staging / marts / sandbox), functional + access roles, service accounts, a warehouse."),
    (1, "Applies", "default grants, PII tags, and naming standards automatically."),
    (0, "Golden-path templates:", "starter dbt project, connector config, and CI workflow scaffolded for the team."),
    (0, "PR-driven:", "a domain is onboarded by opening a PR — reviewed, planned, applied by CI. No tickets."),
    (0, "Outcome:", "consistent and governed by construction; hours, not weeks; the platform team doesn't become the bottleneck."),
], 0.7, 1.5, 12.0, 4.7, base=18)
takeaway(s, "The paved path makes the governed way the easiest way — self-service without losing control.")
footer(s, 9)
notes(s, "This is the 'factory' payoff and a differentiator vs teams that hand-craft each "
         "domain. A single opinionated Terraform module encodes all standards, so every "
         "new domain is born compliant. Teams self-serve via PR; the platform team reviews "
         "the module, not every request. This is how RBAC actually scales to 10 domains — "
         "it's generated, not hand-written.")

# ============================================================================= SLIDE 10 — RBAC diagram
s = content("RBAC at scale: functional + access roles", "Trust & Governance")

def node(x, y, w, h, txt, fill, tc=WHITE, fs=11, bold=True):
    b = rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text_in(b, txt, size=fs, color=tc, bold=bold); return b

def hint(x, y, w, txt):
    tb, tf = textbox(s, x, y, w, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run()
    r.text = txt.upper(); r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = BLUE; r.font.name = FONT

ytop = 1.95
hint(0.6, ytop-0.35, 2.6, "Users")
node(0.6, ytop+0.15, 2.6, 0.6, "OPENGOV_ANALYST", BLUE, fs=10)
node(0.6, ytop+0.95, 2.6, 0.6, "OPENGOV_DBT_SVC", NAVY2, fs=10)
node(0.6, ytop+1.75, 2.6, 0.6, "human  vs  service", GRAY, fs=10, bold=False)

hint(3.75, ytop-0.35, 2.9, "Functional roles (job)")
node(3.75, ytop+0.05, 2.9, 0.6, "REVOPS_ANALYST", BLUE, fs=11)
node(3.75, ytop+0.85, 2.9, 0.6, "REVOPS_DEVELOPER", BLUE, fs=11)
node(3.75, ytop+1.65, 2.9, 0.6, "REVOPS_ADMIN", NAVY, fs=11)

hint(7.25, ytop-0.35, 2.7, "Access roles (grants)")
node(7.25, ytop+0.05, 2.7, 0.55, "AR_MARTS_R", GREEN, fs=10)
node(7.25, ytop+0.75, 2.7, 0.55, "AR_STAGING_R", GREEN, fs=10)
node(7.25, ytop+1.45, 2.7, 0.55, "AR_RAW_RW", GREEN, fs=10)
node(7.25, ytop+2.15, 2.7, 0.55, "AR_SANDBOX_RW", GREEN, fs=10)

hint(10.5, ytop-0.35, 2.3, "Objects")
node(10.5, ytop+0.05, 2.3, 0.55, "MARTS schema", NAVY2, fs=10)
node(10.5, ytop+0.75, 2.3, 0.55, "STAGING schema", NAVY2, fs=10)
node(10.5, ytop+1.45, 2.3, 0.55, "RAW schema", NAVY2, fs=10)
node(10.5, ytop+2.15, 2.3, 0.55, "SANDBOX schema", NAVY2, fs=10)

for x in (3.3, 6.75, 10.05):
    rect(s, x, ytop+1.1, 0.32, 0.28, ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)

tb, tf = textbox(s, 0.6, ytop+2.75, 12.1, 0.5)
p = tf.paragraphs[0]; r = p.add_run()
r.text = "granted to  ►  users     |     inherit  ►  access roles     |     hold privileges on  ►  objects"
r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GRAY; r.font.name = FONT
takeaway(s, "Access roles are reused building blocks; add a domain = generate a new set — grants never sprawl.  Row-access policies handle multi-tenant isolation.")
footer(s, 10)
notes(s, "Snowflake's recommended two-tier model: access roles hold object privileges; "
         "functional roles (job-based) inherit access roles; users get functional roles. "
         "Why it scales: access roles are reusable primitives, so a 10th domain is a "
         "generated set of roles — no bespoke grant spaghetti. Service accounts get their "
         "own functional roles (dbt vs human analyst). Row-access policies enforce "
         "tenant/region isolation on shared tables. This exact model is what I build in Part 2.")

# (SLIDE removed — generic PII/masking; built concretely in Part 2 P2.3
#  "ARR/amount masking — who sees what" + appendix A1)

# ============================================================================= SLIDE 12 — Observability & SLAs
s = content("Observability — what we built", "Trust & Governance")
bullets(s, [
    (0, "Load observability — PAGE_VIEWS_LOAD_LOG:", "one row per ingestion run (file, records loaded, quarantined, timestamp) — an audit trail that outlives COPY's 64-day history."),
    (0, "Quarantine as a signal:", "bad rows land in PAGE_VIEWS_QUARANTINE, countable per run (PROD: 9 loaded / 5 quarantined) — a reject-rate you can alert on."),
    (0, "Quality as code, gated in CI:", "dbt not_null / unique / accepted_values / relationships run on every preprod & prod build — a red test blocks the deploy."),
    (0, "Native audit surfaces(SQL):", "COPY_HISTORY (14d) for loads, ACCESS_HISTORY (365d) for 'who read ACCOUNT.ARR, when'."),
    (0, "Cost attribution:", "per-role warehouses (OG_<ENV>_<ROLE>_WH) + auto-suspend — spend is attributable to reader / analyst / dbt / ingestion, not one shared blob."),
    (0, "Next step:", "dbt source freshness + a monitoring model over _LOAD_LOG turn these into published freshness/quality SLAs."),
], 0.7, 1.5, 12.0, 4.7, base=16)
takeaway(s, "We built the observable primitives — a per-file load log, a quarantine table, CI-gated tests, and native audit views — so 'did the load run and was it clean?' is a query, not a guess.")
footer(s)
notes(s, "REWRITTEN to what we actually built, not generic SLA talk. Concrete artifacts: the "
         "PAGE_VIEWS_LOAD_LOG row-per-run + quarantine table give load observability and a "
         "reject rate; dbt tests gate every CI build; COPY_HISTORY / ACCESS_HISTORY / "
         "POLICY_REFERENCES are the native audit surfaces (ACCESS_HISTORY answers the ARR "
         "probe); per-role warehouses make cost attributable. Full published SLAs (freshness "
         "targets, alerting to Slack/on-call) are the honest next step on top of these.")

# (SLIDE removed — generic CI/CD; built concretely in Part 2 P2.6 (env promotion)
#  + appendix A3 (CI/CD security & idempotency))

# ============================================================================= SLIDE 14 — Differentiator
s = content("Beyond the brief: making the platform AI-ready", "Differentiating Perspective")
bullets(s, [
    (0, "The risk:", "AI/LLM consumers will happily query ungoverned data and return confident, wrong answers."),
    (0, "Governed semantic / metrics layer:", "define metrics once (revenue, pipeline, churn) → one source of truth for BI, notebooks, and LLMs."),
    (0, "Data contracts:", "producers commit to schema + semantics; breaking changes are blocked in CI — trust shifts left."),
    (0, "AI consumes products, not raw:", "RAG and agents read documented, contracted marts + the semantic layer — never RAW."),
    (0, "Governance already travels — true in this build:", "our tag-masking + RBAC bind to the SESSION, not the tool. An LLM service account holding REVOPS_READER gets NULL for ARR today, exactly like a human reader — there is no AI bypass to build."),
], 0.7, 1.5, 12.0, 4.6, base=17)
takeaway(s, "Trustworthy AI and trustworthy analytics come from the same governed foundation — the semantic layer is the interface.")
footer(s, 14)
notes(s, "This is my point of view. The trap everyone is walking into: pointing LLMs at raw "
         "warehouses. Without a semantic layer, the model guesses join logic and metric "
         "definitions — confidently wrong. Solution: a governed metrics/semantic layer as "
         "the single interface for humans AND AI, plus data contracts enforced in CI so "
         "definitions can't silently drift. Crucially, Snowflake's masking/row policies "
         "apply to the LLM's service account — governance isn't bypassed by AI. This is how "
         "OpenGov gets trustworthy AI, not just more AI.")

# (SLIDE removed — Part-1 recap/transition; the Part 2 divider handles the handoff)

# ============================================================================= shared helpers for Part 2 / Appendix
def divider(title, sub):
    d = blank(); set_bg(d, NAVY)
    rect(d, 0, 0, 0.28, SH.inches, ACCENT)
    place_logo(d, 11.9, 1.0, 0.7)
    tb, tf = textbox(d, 1.1, 2.7, 11.2, 2.4)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = tidy(title)
    r.font.size = Pt(38); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    p2 = tf.add_paragraph(); p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = tidy(sub)
    r2.font.size = Pt(17); r2.font.color.rgb = LAV; r2.font.name = FONT
    return d

def bx(sl, x, y, w, h, txt, fill, tc=WHITE, fs=10.5, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    b = rect(sl, x, y, w, h, fill, shape=shape); text_in(b, txt, size=fs, color=tc, bold=bold); return b

def ar(sl, x, y, w=0.3, h=0.24, d=MSO_SHAPE.RIGHT_ARROW):
    rect(sl, x, y, w, h, ACCENT, shape=d)

def lab(sl, x, y, w, txt, color=BLUE, fs=10):
    tb, tf = textbox(sl, x, y, w, 0.3)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; r = p.add_run()
    r.text = txt.upper(); r.font.size = Pt(fs); r.font.bold = True; r.font.color.rgb = color; r.font.name = FONT

# ============================================================================= PART 2 DIVIDER
divider("Part 2: Hands-On Build",
        "RevOps, built and running. RBAC and masking, ADLS ingestion, dbt marts, and CI/CD, all live in Snowflake.")

# ============================================================================= P2.1 — what's live
s = content("RevOps onboarded — the four deliverables, live", "Hands-On Build")
bullets(s, [
    (0, "RBAC & security:", "6 functional + per-schema access roles, per-developer composite roles, tag-based ARR/amount masking — Terraform, config-driven."),
    (0, "Ingestion:", "ingest_page_views.py — ADLS → RAW with validate / quarantine / load-log / idempotency / backfill (the intentional null user_id is quarantined)."),
    (0, "dbt:", "incremental+deduped staging → RevOps pipeline mart; tests, clean_string macro, mesh public model — native dbt on Snowflake."),
    (0, "CI/CD:", "infra (Terraform, remote state) and dbt both deploy through GitHub Actions to DEV & PROD — nothing applied by hand."),
    (0, "Two live databases:", "OG_DEV_DB (dev sandboxes + preprod) and OG_PROD_DB — both seeded, governed, and built end-to-end."),
], 0.7, 1.5, 12.0, 4.6, base=17)
takeaway(s, "Everything in Part 1's architecture is running in the demo account — provisioned and promoted entirely through code.")
footer(s)
notes(s, "Orient the panel: this is the same architecture, now real. Each bullet is a "
         "deliverable I can open live. Emphasize config-driven + CI-deployed — no click-ops.")

# ============================================================================= P2.2 — RBAC + ingestion-layer access (diagram)
s = content("RBAC as built — who can touch each layer", "Hands-On Build")
top = 1.7
# layer column (objects)
lab(s, 0.55, top-0.32, 3.1, "Layer / schema")
layers = [("SALESFORCE_RAW_FIVETRAN", NAVY2), ("PRODUCT_EVENTS_RAW_ADLS", NAVY2),
          ("STAGING", BLUE), ("MARTS_REVOPS", BLUE), ("MARTS_<domain>", GRAY)]
for i,(t,c) in enumerate(layers):
    bx(s, 0.55, top+0.1+i*0.85, 3.1, 0.62, t, c, fs=10)
# write column
lab(s, 4.35, top-0.32, 4.3, "Write access")
writes = ["REVOPS_INGESTION_FIVETRAN  (owns DDL)",
          "REVOPS_INGESTION_ADLS  (DML only — contract is deployer-owned)",
          "REVOPS_DEVELOPER  (dbt CI, all envs)",
          "REVOPS_DEVELOPER  (dbt CI, all envs)",
          "domain developer role (dbt CI)"]
for i,t in enumerate(writes):
    bx(s, 4.35, top+0.1+i*0.85, 4.3, 0.62, t, GREEN, fs=9.5)
# read column
lab(s, 9.0, top-0.32, 3.75, "Read access")
reads = ["REVOPS_READER · ADMIN",
         "REVOPS_READER · ADMIN",
         "REVOPS_READER · DEVELOPER · ADMIN",
         "REVOPS_ANALYST · READER · ADMIN",
         "domain analyst (that domain's marts only)"]
for i,t in enumerate(reads):
    bx(s, 9.0, top+0.1+i*0.85, 3.75, 0.62, t, LIGHT, tc=INK, fs=9.5)
takeaway(s, "Access roles hold object grants; functional roles compose them; humans hold a composite/functional role — writes to shared+PROD happen only via the dbt CI service role.")
footer(s)
notes(s, "Walk each layer: ingestion roles are scoped to exactly their landing schema "
         "(Fivetran owns its DDL; the ADLS loader is DML-only against a platform-owned "
         "contract). Shared STAGING/MARTS are written ONLY by REVOPS_DEVELOPER through the "
         "dbt CI job — no human holds that role. Readers read all; analysts read marts; a "
         "domain analyst is scoped to that domain's marts (FINANCE_ANALYST in the appendix).")

# ============================================================================= P2.3 — masking as built
s = content("ARR / amount masking — who sees what", "Hands-On Build")
bullets(s, [
    (0, "Tag-based:", "GOVERNANCE.PII_FINANCIAL tag → MASK_PII_FINANCIAL_NUMBER policy; classify ACCOUNT.ARR + OPPORTUNITY.AMOUNT once, protection follows the tag."),
    (0, "Exempt (see real values):", "REVOPS_ADMIN, REVOPS_DEVELOPER, REVOPS_ANALYST — NOT REVOPS_READER."),
    (0, "So, in practice:", ""),
    (1, "Analyst", "— sees amount in the MART (exempt); financials are the point of their job."),
    (1, "Reader / a developer in their dev sandbox", "— sees NULL in RAW/STAGING (readers aren't exempt)."),
    (0, "The operational gotcha:", "the tag→policy binding is not a Terraform resource — every apply drops it, so apply_pii_tags.py re-binds it right after apply (automated in CI)."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "Classify a column, not a table — masking scales to every domain through tags, and CI re-establishes the binding on every deploy.")
footer(s)
notes(s, "Demo-able: query ACCOUNT.ARR as REVOPS_READER (NULL) vs the mart as REVOPS_ANALYST "
         "(real). Stress the tag→policy re-bind is the one non-declarative step, and CI owns "
         "it. Audit via ACCESS_HISTORY answers 'who read ARR, when'.")

# ============================================================================= P2.4 — ingestion as built
s = content("Telemetry ingestion — ADLS → RAW, built to the brief", "Hands-On Build")
bullets(s, [
    (0, "Source:", "hourly JSON in ADLS (og-telemetry/<env>/product_events/page_views/dt=/hr=/) — read via a keyless external stage (storage integration, no secrets)."),
    (0, "Validate & quarantine:", "rows missing event_id / account_id / event_timestamp → PAGE_VIEWS_QUARANTINE; the good rows land in append-only RAW."),
    (0, "Load log:", "one PAGE_VIEWS_LOAD_LOG row per file — file, records loaded, quarantined, timestamp — an audit trail that outlives COPY history."),
    (0, "Idempotent + backfillable:", "per-file dedup on re-run; --path/--backfill for a folder or the full set; event_id dedup deferred to dbt staging."),
    (0, "As run:", "PROD load = 9 rows loaded, 5 quarantined; staging dedup collapsed to 5 unique events."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "RAW is immutable truth: validate at the door, quarantine the bad, log every load — and make re-runs safe.")
footer(s)
notes(s, "This is the brief's ingestion task, 1:1 — only S3 became ADLS (owner decision, "
         "boto3 pattern ported to the Azure Blob SDK). The null-user_id record from the "
         "sample payload is exactly what the quarantine catches. Partial-batch: good rows "
         "commit, bad rows quarantine, never fail the whole file. Lambda-on-arrival and a "
         "second event type (config-driven) are the extension answers.")

# ============================================================================= P2.5 — dbt as built
s = content("dbt — staging → RevOps mart, native on Snowflake", "Hands-On Build")
bullets(s, [
    (0, "Staging (incremental, merge):", "typed + snake_case; CDC (_fivetran_synced / _loaded_at); latest-per-PK via QUALIFY; clean_string on stage_name; soft-deletes dropped."),
    (0, "Mart — mart_revops__pipeline:", "opportunities × accounts + days_to_close, pipeline_stage_bucket, weighted_amount. Public mesh model a spoke consumes."),
    (0, "Tests & contracts:", "not_null + unique PKs, accepted_values on stage_name, relationships opp→account, sources.yml — all green."),
    (0, "Three targets:", "dev → your sandbox (schema__model) · preprod → OG_DEV_DB real schemas · prod → OG_PROD_DB (alias naming)."),
    (0, "Governance travels:", "an apply_column_tags post-hook carries the PII_FINANCIAL tag onto built columns in preprod/prod."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "One project, layered + contracted — incremental where volume demands it, mesh-public where consumers depend on it.")
footer(s)
notes(s, "Native dbt on Snowflake: the repo runs as a DBT PROJECT object, executed in-account "
         "as REVOPS_DEVELOPER — the CI runner never builds models. Incremental vs full-refresh "
         "is the volume decision (appendix). mart_revops__pipeline is the mesh public model.")

# ============================================================================= P2.6 — CI/CD env promotion (diagram)
s = content("CI/CD — how code becomes DEV, QA & PROD", "Hands-On Build")
# dbt lane
lab(s, 0.55, 1.45, 2.2, "dbt (models)", ACCENT, fs=11)
bx(s, 0.55, 1.8, 2.55, 0.95, "feature/<you>\n→ target dev", NAVY2, fs=10, bold=True)
ar(s, 3.2, 2.15)
bx(s, 3.6, 1.8, 3.0, 0.95, "your SANDBOX\nREVOPS_DEV_<NAME>\n(OG_DEV_DB)", GREEN, fs=9.5, bold=True)
ar(s, 6.7, 2.15)
bx(s, 7.1, 1.8, 2.6, 0.95, "push dev\n→ target preprod\n(OG_DEV_DB real schemas)", BLUE, fs=9, bold=True)
ar(s, 9.8, 2.15)
bx(s, 10.2, 1.8, 2.55, 0.95, "push main\n→ target prod\n(OG_PROD_DB)", NAVY, fs=9.5, bold=True)
# infra lane
lab(s, 0.55, 3.15, 2.2, "infra (RBAC/gov)", ACCENT, fs=11)
bx(s, 0.55, 3.5, 3.0, 0.95, "PR to main\n→ plan + drift check\n(posted as PR comment)", NAVY2, fs=9.5, bold=True)
ar(s, 3.65, 3.85)
bx(s, 4.05, 3.5, 3.4, 0.95, "merge main\n→ terraform apply\n+ apply_pii_tags (per env)", BLUE, fs=9.5, bold=True)
ar(s, 7.55, 3.85)
bx(s, 7.95, 3.5, 4.8, 0.95, "Snowflake OG_DEV_DB + OG_PROD_DB\nroles · grants · masking · stages", NAVY, fs=9.5, bold=True)
# env legend
band = rect(s, 0.55, 4.95, 12.2, 1.0, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE); _noline(band)
tb, tf = textbox(s, 0.8, 5.05, 11.8, 0.85)
for i,(k,v) in enumerate([
    ("DEV", "per-developer sandboxes in OG_DEV_DB; build and iterate freely"),
    ("QA / preprod", "real STAGING/MARTS in OG_DEV_DB; integration build on merge to dev"),
    ("PROD", "OG_PROD_DB; built on merge to main, protected environment + approval")]):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.line_spacing=1.05
    r0=p.add_run(); r0.text=f"{k}:  "; r0.font.bold=True; r0.font.size=Pt(11); r0.font.color.rgb=BLUE; r0.font.name=FONT
    r1=p.add_run(); r1.text=v; r1.font.size=Pt(11); r1.font.color.rgb=INK; r1.font.name=FONT
footer(s)
notes(s, "The env story: DEV database hosts BOTH the per-developer sandboxes (target dev) and "
         "the QA/preprod integration schemas (target preprod, built on merge to dev). PROD is "
         "a separate database, built on merge to main behind a protected GitHub Environment. "
         "Auth is key-pair from secrets; secrets are only available post-merge, so a malicious "
         "PR can't reach them. A manual full-refresh workflow handles incremental logic fixes.")

# ============================================================================= P2.7 — developer isolation
s = content("Developer self-service, safely isolated", "Hands-On Build")
bullets(s, [
    (0, "One composite role per developer:", "DEV_<NAME> = REVOPS_READER (read all) + write to only their own sandbox REVOPS_DEV_<NAME>."),
    (0, "Single active role:", "it's their default role and secondary roles are OFF — every session is exactly one role, so masking checks are unambiguous."),
    (0, "Two developers never collide:", "separate sandboxes, <schema>__<model> naming; can't touch shared STAGING/MARTS or anything in PROD."),
    (0, "Local or Snowsight:", "dbt-core locally (username+password, build --target dev) or a Snowsight Workspace on a feature branch."),
    (0, "Onboard = two CSV rows:", "add to environments.csv developers + human_users.csv → composite role + sandbox generated by Terraform in a PR."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "Developers move fast in their own sandbox; the platform guarantees they can't break shared or production data.")
footer(s)
notes(s, "This is the 'paved path' for a person. The composite role is the key design: read "
         "everything, write only your sandbox, one role active — no secondary-role ambiguity. "
         "Onboarding is a two-line PR (demonstrated live with AKASHPAHILWAN, SOURABH_SHINDE, "
         "ANUJKUMAR).")

# ============================================================================= FUTURE ENHANCEMENTS DIVIDER
divider("Future Enhancements", "Where the platform goes next. Each is an additive step on the foundation, not a rebuild.")

# ============================================================================= FE.1 — ingestion / dbt / governance
s = content("Future enhancements: ingestion, dbt & governance", "Roadmap")
bullets(s, [
    (0, "Ingestion:", ""),
    (1, "Snowpipe auto-ingest", "— event-driven per-file load; retires the scheduled COPY + watermark (Event Grid / Lambda on arrival)."),
    (1, "Config-driven multi-event framework", "— onboard a second event type with config, not copy-pasted code."),
    (0, "dbt / transforms:", ""),
    (1, "Slim CI (state:modified)", "— build only what changed; contracts on staging models; exposures for BI/AI lineage."),
    (1, "Domain spoke repos", "— consume the hub's public mart cross-project (full dbt Mesh)."),
    (0, "Governance & security:", ""),
    (1, "Row-access policies", "— tenant / region isolation on shared tables."),
    (1, "OIDC federated CI auth", "— retire long-lived key-pairs; env-specific DEV_* / PROD_* role split for per-env isolation."),
], 0.7, 1.45, 12.0, 4.85, base=15.5)
takeaway(s, "The foundation is built; each of these is an additive step on top of it, not a rebuild.")
footer(s)
notes(s, "Roadmap, part 1. Snowpipe replaces the scheduled COPY once files land continuously "
         "(the ingestion design already keeps the same append-only contract). Slim CI + "
         "staging contracts + spoke repos complete the dbt Mesh story. Row-access policies "
         "and OIDC are the security hardening; env-specific DEV_*/PROD_* roles (SnowOps-style) "
         "give per-environment isolation beyond the account-wide functional roles.")

# ============================================================================= FE.2 — observability / cost / AI
s = content("Future enhancements: observability, cost & AI", "Roadmap")
bullets(s, [
    (0, "Observability & SLAs:", ""),
    (1, "Published data SLAs", "— dbt source freshness + a monitoring model over _LOAD_LOG; freshness / quality targets we measure."),
    (1, "Alerting", "— test & freshness breaches routed to Slack / on-call."),
    (0, "Cost / FinOps:", ""),
    (1, "Per-domain cost dashboards", "— resource monitors + query budgets on the per-role warehouses."),
    (0, "AI-readiness:", ""),
    (1, "Governed semantic / metrics layer", "— one metric definition for BI, notebooks and LLMs; no more confidently-wrong AI."),
    (1, "AI reads products, not raw", "— RAG / agents on contracted marts; masking already applies to an LLM service account."),
], 0.7, 1.45, 12.0, 4.85, base=15.5)
takeaway(s, "Trustworthy AI and trustworthy analytics come from the same governed foundation; the semantic layer is the next interface.")
footer(s)
notes(s, "Roadmap, part 2. Observability grows from the primitives we built (_LOAD_LOG, dbt "
         "tests) into published SLAs + alerting. Cost dashboards use the per-role warehouse "
         "split for attribution. The AI-readiness bets are the differentiator made concrete: "
         "a governed semantic layer as the single interface, contracts across domains, and "
         "governance that already travels to any service account including an LLM's.")

# ============================================================================= FINAL CLOSE
s = content("Close: one governed foundation, analytics & AI on top", "Wrap-up")
bullets(s, [
    (0, "Foundation → pipelines → self-service → trust:", "one IaC-first, governed platform — and it's running, not just drawn."),
    (0, "Buy the commodity, build the differentiator:", "Fivetran for SaaS; custom Python for telemetry; native dbt for transforms."),
    (0, "Guardrails as paved paths:", "config-driven RBAC + tag masking, per-developer sandboxes, CI/CD to DEV & PROD."),
    (0, "Scales by construction:", "a 10th domain is the same CSV rows + module — roles, masking, and CI all generated."),
], 0.7, 1.5, 12.0, 4.6, base=17)
takeaway(s, "The measure of the platform: how fast a federated team ships something trustworthy on top of it — safely, and on their own.")
footer(s)
notes(s, "Land the thesis: platform-as-product, governed by default, self-service, running "
         "end-to-end. Invite deep-dive questions — the appendix has the internals.")

# ============================================================================= APPENDIX DIVIDER
divider("Technical Appendix", "Deep dives for the panel's probes. Pull these up on demand.")

# ============================================================================= A1 — RBAC scale / domain-scoped analyst
s = content("Scaling RBAC to 10 domains — least privilege by construction", "Appendix · Governance")
bullets(s, [
    (0, "Access roles are reusable primitives:", "AR_<ENV>_<SCHEMA>_R/W hold the grants; a new domain = a generated set, never bespoke grant spaghetti."),
    (0, "Domain-scoped analysts:", "REVOPS_ANALYST reads MARTS_* (all marts); FINANCE_ANALYST reads MARTS_FINANCE ONLY — an exact-schema grant, not the wildcard."),
    (0, "Proven live:", "FINANCE_ANALYST sees only MARTS_FINANCE and is DENIED MARTS_REVOPS; the RevOps analyst still sees both."),
    (0, "Service vs human:", "REVOPS_DEVELOPER (write, incl PROD) is service-only — sync_config.py rejects granting it to a human; humans get read + their own sandbox."),
    (0, "Audit:", "ACCESS_HISTORY answers 'who read ACCOUNT.ARR and when'; POLICY_REFERENCES proves the mask was attached at read time."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "Least privilege scales because it's generated from config — each domain analyst is confined to that domain's marts by an exact grant.")
footer(s)
notes(s, "Direct answer to 'how does this scale to 10 domains + how do you handle service vs "
         "human + audit'. The FINANCE_ANALYST demo is the concrete proof of domain-scoped "
         "least privilege — built and verified live.")

# ============================================================================= A2 — incremental vs full-refresh
s = content("Incremental vs full-refresh — and fixing incrementals", "Appendix · dbt")
bullets(s, [
    (0, "Incremental (merge):", "high-volume/append sources (telemetry, CDC tables) — build only rows since the last watermark; cheap, fast."),
    (0, "Full-refresh:", "small dimensions, or after a logic change — rebuild from scratch."),
    (0, "The trap:", "a normal build only merges NEW rows, so a logic fix to an incremental model never re-cleans already-materialized rows."),
    (0, "The fix, as a button:", "a manual full-refresh workflow — branch (dev→preprod / main→prod) × scope (all | only models changed in the last commit)."),
    (0, "Real example:", "the page_views user_id VARIANT-null fix needed a full-refresh to purge 2 historical bad rows from PROD."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "Pick incremental for volume, full-refresh for correctness — and make the full-refresh a one-click, scoped workflow.")
footer(s)
notes(s, "The 'incremental vs full-refresh' probe. The non-obvious operational point: an "
         "incremental logic fix requires a full-refresh, which is why there's a dedicated "
         "workflow with all/modified scope. Grounded in a real bug we fixed.")

# ============================================================================= A3 — CI/CD security & idempotency
s = content("CI/CD security, idempotency & standards", "Appendix · Delivery")
bullets(s, [
    (0, "No static secrets:", "key-pair (JWT) for Snowflake, storage key for the TF backend — all GitHub secrets; the private key touches only a temp file in the job."),
    (0, "Malicious-PR safe:", "secrets are exposed only to post-merge apply jobs; PR-triggered plan never sees them; PROD gated by a protected Environment + approval."),
    (0, "Remote state:", "Terraform state in an azurerm backend (ADLS), so CI and local share one source of truth — no drift, no local state to lose."),
    (0, "Idempotent & half-failure safe:", "config→JSON→for_each reruns cleanly; a drift check fails the PR if CSVs and generated JSON disagree; apply_pii_tags skips absent objects."),
    (0, "Standards enforced:", "sync_config.py validates every reference (unknown role/schema/env) and rejects a service-only role reaching a human — the linter is the reviewer."),
], 0.7, 1.5, 12.0, 4.6, base=16.5)
takeaway(s, "The pipeline is the control plane: reviewed, idempotent, secret-safe, and standards-enforcing — provisioning you can trust.")
footer(s)
notes(s, "Hits every CI/CD probe: OIDC/key-pair, protecting secrets from fork PRs, "
         "idempotency + half-failure safety, linting/naming. The drift check + sync_config "
         "validation are our 'linting'; remote state + post-merge-only secrets are the "
         "security spine.")

# ============================================================================= SLIDE — BONUS (appendix): VARIANT internals
s = content("Why VARIANT stays fast at scale", "Appendix · Internals")
bullets(s, [
    (0, "Not a text blob:", "each micro-partition shreds consistent JSON paths into typed, compressed sub-columns with min/max stats — native-column treatment."),
    (0, "Projection works:", "selecting payload:properties.duration_ms reads only that sub-column off disk — never the whole document."),
    (0, "Pruning works:", "filters on extracted paths skip whole micro-partitions via metadata — scan cost tracks the answer, not the table."),
    (0, "When it degrades:", "wildly heterogeneous keys defeat sub-columnarization — promote hot filter / dedup keys to real columns (event_id)."),
    (0, "Discipline & limits:", "16 MB per value (events are KBs); missing key → NULL, never an error; TRY_CAST + tests absorb type drift; marts never read payload."),
], 0.65, 1.5, 7.05, 4.7, base=16.5)
code_box(s, 8.0, 1.6, 4.78, 4.55, [
    "-- inside one micro-partition,",
    "-- payload VARIANT is stored as:",
    "--  :event_id            -> sub-col",
    "--  :account_id          -> sub-col",
    "--  :properties.duration_ms -> sub-col",
    "--  ...each with min/max stats",
    "",
    "-- projection: reads ONE sub-column",
    "SELECT payload:properties.duration_ms",
    "FROM raw.product_events.page_views",
    "",
    "-- pruning: skips micro-partitions",
    "WHERE payload:event_timestamp",
    "      ::timestamp_ntz >= '2026-07-01'",
    "",
    "-- drift-safe typing in staging",
    "TRY_CAST(payload:properties.duration_ms",
    "         AS INT)  -- bad value -> NULL",
], title="VARIANT internals (sketch)")
takeaway(s, "Schema-on-read is nearly free at query time — consistent paths are columnarized on write; type once in staging, downstream never parses JSON.")
footer(s)
notes(s, "BONUS / HIDDEN — pull up only if the panel digs into VARIANT performance. Key "
         "message: VARIANT is not a JSON string column. On write, Snowflake examines paths "
         "that occur consistently across rows in a micro-partition and stores each as its "
         "own typed, compressed sub-column with min/max metadata — so path projection and "
         "partition pruning behave like native columns. It degrades only when keys are "
         "wildly inconsistent row-to-row; the fix is promoting hot keys (event_id) to real "
         "columns, which we do anyway for dedup. Drift behavior: missing key reads as NULL "
         "(never an error), new keys are invisible until a PR selects them, and TRY_CAST + "
         "dbt tests turn type drift into a quarantine signal instead of a failed job. "
         "16 MB/value limit is ~4 orders of magnitude above a telemetry event. One-breath "
         "close: 'stable RAW DDL, columnarized paths, type once in incremental staging — "
         "downstream never touches JSON.'")
# appendix slide — visible in slideshow

# ============================================================================= SLIDE 17 — BONUS (hidden): S3 layout & hourly loads
s = content("ADLS layout & hourly loads — append-only RAW", "Appendix · Internals")
bullets(s, [
    (0, "One prefix per source object:", "source/object/dt=YYYY-MM-DD/hr=HH/ — hourly runs list one narrow prefix; backfill = re-point at a date range; lifecycle rules attach per prefix."),
    (0, "Append new files, never replace:", "RAW is an immutable log, not a mirror — 'latest' applies to which files we read this run, not to what we keep."),
    (0, "Idempotent in three layers:", "COPY's 64-day file-load history (re-run loads zero) · _LOAD_LOG watermark + trailing-window rescan for late files · event_id dedup in staging."),
    (0, "Full refresh without mutation:", "COPY FORCE=TRUE into a fresh table, then ALTER TABLE SWAP — rebuild history, never edit it."),
    (0, "COPY now, Snowpipe next:", "scheduled Python + COPY keeps validation / quarantine / _LOAD_LOG in code; production evolution: Snowpipe auto-ingest per S3 event, no scheduler."),
], 0.65, 1.5, 7.05, 4.7, base=16.5)
code_box(s, 8.0, 1.6, 4.78, 4.55, [
    "azure://snowopssa.blob.core.windows.net/",
    "  og-telemetry/prod/product_events/",
    "    page_views/dt=2026-07-12/hr=14/",
    "      page_views_2026-07-12_14.json",
    "    .../hr=15/ ...",
    "",
    "-- hourly job: append new files only",
    "COPY INTO product_events_raw_adls.page_views",
    "FROM @page_views_stage/dt=2026-07-12/",
    "FILE_FORMAT = ff_json;",
    "-- 64-day file history: re-run = 0 rows",
    "",
    "-- staging: producer resends -> 1 row",
    "QUALIFY ROW_NUMBER() OVER (",
    "  PARTITION BY event_id",
    "  ORDER BY loaded_at DESC) = 1",
], title="land hourly, append-only")
takeaway(s, "Copy only the latest files, keep every row — append-only RAW plus three idempotency layers makes any re-run or backfill safe.")
footer(s)
notes(s, "BONUS / HIDDEN — pull up if asked about bucket layout, hourly cadence, or 'why not "
         "Snowpipe'. Hive-style dt=/hr= prefixes sort lexicographically by time, so the "
         "watermark is just the last processed prefix; re-scan a trailing 24-48h window so "
         "late-arriving files in an old hour still load (COPY's file history dedups the "
         "overlap for free). Append-only stance: never 'replace with latest' — RAW keeps "
         "every loaded row; dedup and current-state are computed in staging (QUALIFY on "
         "event_id). Snowpipe question: for the case-study deliverable I keep a scheduled "
         "Python job + COPY, because the assessed logic — schema validation, quarantine, "
         "_LOAD_LOG summary — lives in that code, and hourly batch on an XS warehouse is "
         "cheap. In production, with files landing continuously, I'd flip to Snowpipe: S3 "
         "event notification -> per-file auto-ingest, exactly-once per file, ~1-minute "
         "latency, no scheduler or watermark logic at all; validation then moves to a "
         "post-load stream/task. Same table, same append-only contract — only the trigger "
         "changes. 64-day detail: COPY's per-file load history EXPIRES after 64 days. A "
         "backfill touching older files is silently SKIPPED by default (status unknown -> "
         "skip, no error — a silent gap, not duplicates). Deliberate old backfill: set "
         "LOAD_UNCERTAIN_FILES=TRUE (staging's event_id dedup absorbs any repeats) or "
         "FORCE=TRUE into a fresh table + SWAP. This expiry is exactly why the design has "
         "three layers — COPY history expires, a watermark can miss late files, so event_id "
         "dedup in staging is the backstop that never expires.")
# appendix slide — visible in slideshow

# ============================================================================= SLIDE 18 — BONUS (hidden): COPY logs & 64-day horizon
s = content("COPY INTO — load logs & the 64-day horizon", "Appendix · Internals")
bullets(s, [
    (0, "COPY returns a per-file result set:", "rows_parsed / rows_loaded / errors per file — the Python job persists it to _LOAD_LOG; ours never expires."),
    (0, "Built-in audit views:", "COPY_HISTORY table function = 14 days, per table; ACCOUNT_USAGE = 365 days, account-wide (~90 min lag)."),
    (0, "The 64-day dedup horizon:", "the internal metadata COPY uses to skip already-loaded files expires after 64 days — separate from the audit views."),
    (0, "Beyond 64 days = silent skip:", "unknown-status files are skipped — a silent gap, not duplicates; backfill via LOAD_UNCERTAIN_FILES=TRUE or FORCE + SWAP."),
    (0, "Why three layers:", "COPY history expires · watermarks miss late files · event_id dedup in staging never expires."),
], 0.65, 1.5, 7.05, 4.7, base=16.5)
code_box(s, 8.0, 1.6, 4.78, 4.55, [
    "-- COPY's return: per-file summary",
    "--  file | status | rows_parsed |",
    "--  rows_loaded | first_error ...",
    "--  -> Python writes it to _LOAD_LOG",
    "",
    "-- audit: last 14 days, this table",
    "SELECT file_name, last_load_time,",
    "       row_count, error_count",
    "FROM TABLE(information_schema",
    "  .copy_history(",
    "   table_name => 'PAGE_VIEWS',",
    "   start_time => dateadd(day, -14,",
    "     current_timestamp())));",
    "",
    "-- audit: 365 days, account-wide",
    "SELECT * FROM snowflake",
    "  .account_usage.copy_history;",
    "",
    "-- backfill older than 64 days",
    "COPY INTO ...",
    "  LOAD_UNCERTAIN_FILES = TRUE;",
], title="every load is observable")
takeaway(s, "Snowflake logs every COPY (14-day / 365-day views) and returns a per-file summary — persisting it in _LOAD_LOG gives an audit trail that outlives every retention window.")
footer(s)
notes(s, "BONUS / HIDDEN — pull up for 'can we see what COPY loaded?' or the 64-day "
         "follow-up. Three observable layers: (1) COPY INTO itself returns a result set — "
         "one row per file with status, rows_parsed, rows_loaded, errors_seen, first_error "
         "— and the Part 2 Python job captures exactly this into _LOAD_LOG, which is the "
         "brief's load-summary requirement and lives forever. (2) COPY_HISTORY / "
         "LOAD_HISTORY: INFORMATION_SCHEMA table functions keep 14 days per table; "
         "SNOWFLAKE.ACCOUNT_USAGE views keep 365 days account-wide with up to ~90 min "
         "latency — these are audit/observability surfaces. (3) Distinct from both: the "
         "internal per-file dedup metadata COPY consults to skip already-loaded files — "
         "64-day lifetime, not directly queryable. Beyond 64 days the load status is "
         "'unknown' and COPY silently SKIPS the file by default (gap, not duplicates); "
         "deliberate backfills set LOAD_UNCERTAIN_FILES=TRUE and let staging's event_id "
         "dedup absorb repeats, or rebuild with FORCE=TRUE + SWAP. Close with the "
         "three-layer line: history expires, watermarks miss late files, staging dedup "
         "never expires.")
# appendix slide — visible in slideshow

# ----------------------------------------------------------------------------- save
out = __file__.rsplit("\\", 1)[0] + "\\OpenGov_Data_Platform_CaseStudy.pptx"
prs.save(out)
print("Saved:", out)
print("Slides:", len(prs.slides._sldIdLst))
